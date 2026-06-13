"""
AI Assistant Development - PowerPoint Presentation Generator
Generates a professional presentation with project details, screenshots, links, and content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (matching app theme - dark blue/purple)
DARK_BG = RGBColor(4, 5, 11)  # #04050b
PRIMARY_COLOR = RGBColor(139, 92, 246)  # Purple
ACCENT_COLOR = RGBColor(59, 130, 246)  # Blue
TEXT_COLOR = RGBColor(255, 255, 255)  # White
MUTED_COLOR = RGBColor(156, 163, 175)  # Gray

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list, add_line=True):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add title underline
    if add_line:
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(2), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(12)

def add_image_slide(prs, title, image_path):
    """Add slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add image if it exists
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    else:
        # Add placeholder text
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = f"[Screenshot: {os.path.basename(image_path)}]"
        placeholder_frame.paragraphs[0].font.size = Pt(20)
        placeholder_frame.paragraphs[0].font.color.rgb = MUTED_COLOR
        placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 1: Title Slide
add_title_slide(prs, 
    "AI Assistant Development",
    "A Web-Based AI Assistant Demonstrating Prompt Engineering\nand Google Gemini API Integration\n\nInternship Prompt Engineering Major Project\nDeveloper: Bugata Pravallika | 2026"
)

# Slide 2: Project Objective
add_content_slide(prs, "Project Objective", [
    "• Educational Purpose: Illustrate prompt optimization and LLM instructions",
    "• Demonstrate Prompt Engineering: Compare multiple prompt templates side-by-side",
    "• Practical Utility: Production-ready assistant for QA, Summarization, Content Generation, and Study Advice",
    "• Deployment Ready: GitHub + Render integration for seamless deployment"
])

# Slide 3: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "• Prompt Sensitivity: LLM outputs vary significantly based on prompt structure and phrasing",
    "• Access Barriers: Users struggle to write optimized prompts",
    "• Interface Gap: No system for testing multiple prompt templates side-by-side",
    "• Evaluation Needs: Lack of feedback mechanism for continuous prompt optimization"
])

# Slide 4: Key Features
add_content_slide(prs, "Key Features", [
    "✓ Question Answering with customizable templates",
    "✓ Text Summarization: Extract key points automatically",
    "✓ Creative Content Generation: Poems, stories, essays",
    "✓ Study Advice & Recommendations for career guidance",
    "✓ Live Statistics Dashboard with real-time metrics"
])

# Slide 5: Prompt Engineering Framework
add_content_slide(prs, "Prompt Engineering Framework", [
    "🎭 Dual Persona System:",
    "   • Friendly Persona: Warm, beginner-friendly explanations",
    "   • Academic Persona: Formal, rigorous, detailed explanations",
    "\n📝 Function Templates:",
    "   • 3 distinct templates per function for comparison",
    "   • Dynamic variable substitution for user inputs"
])

# Slide 6: System Architecture
add_content_slide(prs, "System Architecture", [
    "Frontend → HTML5 + CSS3 (Glassmorphism UI) + JavaScript (ES6+)",
    "API Layer → AJAX requests with JSON payloads",
    "Backend → Flask (Python) server handling requests",
    "AI Integration → Google Gemini API (gemini-1.5-flash)",
    "Response Processing → Markdown parsing with marked.js",
    "Data Logging → Feedback saved to feedback.txt for analytics"
])

# Slide 7: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "Backend: Python 3.12 + Flask",
    "Frontend: HTML5, CSS3 (Vanilla, Flexbox/Grid), JavaScript ES6+",
    "AI Engine: Google Generative AI SDK (Gemini API)",
    "Production: Gunicorn (WSGI server)",
    "Infrastructure: GitHub (Version Control) + Render (Cloud Deployment)",
    "Additional: Marked.js (Markdown parsing), FontAwesome (Icons)"
])

# Slide 8: UI/UX Design
add_content_slide(prs, "UI/UX Design Highlights", [
    "🌙 Dark Mode by Default: Premium deep navy, purple, and blue palette",
    "✨ Glassmorphism: Translucent cards with blur effects and ambient glow",
    "⚡ Interactive Controls: Fluid template picker cards and responsive layout",
    "💬 ChatGPT-Style: Full Markdown support, code syntax highlighting",
    "📊 Live Dashboard: Real-time query counts and feedback metrics"
])

# Slide 9: Feedback Mechanism
add_content_slide(prs, "Feedback Loop & Data Logging", [
    "👍 Immediate User Response: 'Helpful' / 'Not Helpful' options",
    "📁 File-Based Logging: Records timestamps, function type, persona, templates, and ratings",
    "📈 Continuous Improvement: Log data for refining prompts and personas",
    "✨ Thank You Animation: Interactive feedback confirmation"
])

# Slide 10: Screenshots - Hero Section
add_image_slide(prs, "Screenshot: Landing Page", 
    "screenshots/hero-section.png"
)

# Slide 11: Screenshots - Assistant Interface
add_image_slide(prs, "Screenshot: AI Assistant Interface", 
    "screenshots/assistant-interface.png"
)

# Slide 12: Screenshots - Dashboard
add_image_slide(prs, "Screenshot: Statistics Dashboard", 
    "screenshots/dashboard.png"
)

# Slide 13: Project Links & Deployment
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Project Links & Deployment"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

# Add line
line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(3), Inches(0))
line.line.color.rgb = ACCENT_COLOR
line.line.width = Pt(3)

# Links content
links_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
links_frame = links_box.text_frame
links_frame.word_wrap = True

link_items = [
    "🔗 GitHub Repository:",
    "https://github.com/BugataPravallika/ai-assistant-development",
    "",
    "🚀 Live Deployment (Render):",
    "https://ai-assistant-development-zqek.onrender.com",
    "",
    "📋 Features:",
    "✓ Production-ready Flask application",
    "✓ Google Gemini API integration",
    "✓ Responsive dark-mode UI",
    "✓ Real-time feedback analytics"
]

for i, item in enumerate(link_items):
    if i > 0:
        links_frame.add_paragraph()
    p = links_frame.paragraphs[i]
    p.text = item
    if "https://" in item:
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_COLOR
        p.font.bold = True
    else:
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
    p.space_before = Pt(6)

# Slide 14: Conclusion
add_content_slide(prs, "Key Takeaways", [
    "✅ Successfully demonstrated Prompt Engineering principles in a production web application",
    "✅ Implemented dual-persona system for diverse user needs",
    "✅ Built responsive, modern UI with glassmorphic design",
    "✅ Deployed to Render with GitHub integration",
    "✅ Created feedback mechanism for continuous improvement"
])

# Save presentation
output_path = "AI_Assistant_Development_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"\n📊 Total slides: {len(prs.slides)}")
print(f"📁 Location: {os.path.abspath(output_path)}")
